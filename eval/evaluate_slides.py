
import os
import sys
import re
import math
import argparse
import json
import numpy as np
from typing import List, Dict, Any
from pptx import Presentation
import fitz
from collections import Counter

# Try importing specialized AI libs, handle failures gracefully
try:
    import torch
    import torch.nn.functional as F
    from transformers import AutoTokenizer, AutoModel
    HAS_TRANSFORMERS = True
except ImportError:
    HAS_TRANSFORMERS = False

# Reuse existing project logic for section extraction to ensure fair comparison
sys.path.append(os.getcwd())
try:
    from paper2ppt_core.sections import split_into_sections, normalize_heading
    from paper2ppt_core.io import load_input_paper
except ImportError:
    print("[ERROR] Could not import project core modules. Run this from the project root.")
    sys.exit(1)

def simple_tokenize(text):
    return re.findall(r'\w+', text.lower())

# ==========================================
# METRIC: ROUGE (Simplified Implementation)
# ==========================================
def calculate_rouge_l(reference: str, candidate: str) -> float:
    """
    Computes ROUGE-L (Longest Common Subsequence) F1 score.
    """
    ref_tokens = simple_tokenize(reference)
    cand_tokens = simple_tokenize(candidate)
    
    if not ref_tokens or not cand_tokens:
        return 0.0

    # DP for LCS
    m, n = len(ref_tokens), len(cand_tokens)
    dp = [[0] * (n + 1) for _ in range(m + 1)]
    
    for i in range(1, m + 1):
        for j in range(1, n + 1):
            if ref_tokens[i - 1] == cand_tokens[j - 1]:
                dp[i][j] = dp[i - 1][j - 1] + 1
            else:
                dp[i][j] = max(dp[i - 1][j], dp[i][j - 1])
                
    lcs = dp[m][n]
    
    prec = lcs / n if n > 0 else 0
    rec = lcs / m if m > 0 else 0
    
    if prec + rec == 0:
        return 0.0
        
    f1 = 2 * ((prec * rec) / (prec + rec))
    return f1

# ==========================================
# METRIC: SEMANTIC SIMILARITY (Embeddings)
# ==========================================
class SemanticEvaluator:
    def __init__(self):
        self.model = None
        self.tokenizer = None
        if HAS_TRANSFORMERS:
            # unique small model for evaluation
            model_id = "sentence-transformers/all-MiniLM-L6-v2" 
            try:
                print(f"[EVAL] Loading semantic model: {model_id}...")
                self.tokenizer = AutoTokenizer.from_pretrained(model_id)
                self.model = AutoModel.from_pretrained(model_id)
            except Exception as e:
                print(f"[WARN] Failed to load semantic model: {e}. Using TF-IDF fallback.")
                self.model = None

    def get_embedding(self, text: str):
        if not self.model:
            return None
        
        inputs = self.tokenizer(text, return_tensors="pt", padding=True, truncation=True, max_length=512)
        with torch.no_grad():
            outputs = self.model(**inputs)
        
        # Mean Pooling - Take attention mask into account for correct averaging
        token_embeddings = outputs.last_hidden_state
        attention_mask = inputs['attention_mask']
        
        input_mask_expanded = attention_mask.unsqueeze(-1).expand(token_embeddings.size()).float()
        return torch.sum(token_embeddings * input_mask_expanded, 1) / torch.clamp(input_mask_expanded.sum(1), min=1e-9)

    def compute_similarity(self, text1: str, text2: str) -> float:
        if not self.model:
            # Fallback: Jaccard Similarity on tokens
            s1 = set(simple_tokenize(text1))
            s2 = set(simple_tokenize(text2))
            if not s1 or not s2: return 0.0
            return len(s1.intersection(s2)) / len(s1.union(s2))

        emb1 = self.get_embedding(text1)
        emb2 = self.get_embedding(text2)
        
        # Cosine similarity
        cosine_sim = F.cosine_similarity(emb1, emb2)
        return cosine_sim.item()

# ==========================================
# METRIC: READABILITY (Flesch-Kincaid)
# ==========================================
def calculate_readability(text: str) -> float:
    """ Returns Flesch-Kincaid Grade Level (approx) """
    if not text.strip():
        return 0.0
        
    sentences = re.split(r'[.!?]+', text)
    sentences = [s for s in sentences if s.strip()]
    num_sentences = len(sentences) or 1
    
    words = simple_tokenize(text)
    num_words = len(words) or 1
    
    # Syllable approximation
    def count_syllables(word):
        word = word.lower()
        count = 0
        vowels = "aeiouy"
        if word[0] in vowels:
            count += 1
        for i in range(1, len(word)):
            if word[i] in vowels and word[i - 1] not in vowels:
                count += 1
        if word.endswith("e"):
            count -= 1
        if count == 0:
            count += 1
        return count
        
    num_syllables = sum(count_syllables(w) for w in words)
    
    # FKGL Formula
    score = 0.39 * (num_words / num_sentences) + 11.8 * (num_syllables / num_words) - 15.59
    return score

# ==========================================
# MAIN PIPELINE
# ==========================================
def evaluate_pipeline(pdf_path, pptx_path, output_json="eval_report.json"):
    
    print(f"--- Starting Evaluation ---")
    print(f"Input PDF: {pdf_path}")
    print(f"Input PPTX: {pptx_path}")

    # 1. Extract PDF Sections (Ground Truth)
    print("1. Extracting PDF content...")
    pages_text, _ = load_input_paper(pdf_path)
    # Re-use the project's own section splitter to map 1:1 if possible
    raw_sections = split_into_sections(pages_text)
    
    # Map normalized title -> full text
    pdf_map = {}
    for sec in raw_sections:
        # Use simple key for matching
        key = normalize_heading(sec['raw_title'])
        if key == "section": 
            # try to make it unique if generic
            key = sec['raw_title'].lower()
        pdf_map[key] = sec['text']

    # 3. Extract PPTX Content (Generated Hypothesis)
    print("2. Extracting PPTX content...")
    if not os.path.exists(pptx_path):
        print(f"[ERROR] PPTX not found: {pptx_path}")
        return

    prs = Presentation(pptx_path)
    slides_data = []
    
    for i, slide in enumerate(prs.slides):
        title = "Untitled"
        
        # Strategy 1: Standard Placeholder
        if slide.shapes.title:
            title = slide.shapes.title.text
        
        # Strategy 2: Heuristic (Top-most text box)
        if title == "Untitled":
            # Search for text box near the top
            possible_titles = []
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    # Check if it's near the top (e.g., top < 1.0 inch)
                    # PPTX units are EMUs (914400 per inch)
                    # 1 inch = 914400 EMUs. Let's say top < 1.2 inches is a title candidate
                    if shape.top < 1100000: 
                        text = shape.text_frame.text.strip()
                        if text:
                            possible_titles.append((shape.top, text))
            
            # Sort by top position, take the highest one
            if possible_titles:
                possible_titles.sort(key=lambda x: x[0])
                title = possible_titles[0][1]

        # Extract Bullets (content)
        bullets = []
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                text = shape.text_frame.text.strip()
                # Skip if it's the title we just found
                if text == title.strip():
                    continue
                # Skip valid footers/metadata
                if "Auto-generated" in text or "Questions?" in text:
                    continue
                
                bullets.append(text)
        
        slide_text = " ".join(bullets)
        
        # Extract Notes (Narration)
        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            
        slides_data.append({
            "id": i+1,
            "title": title,
            "text": slide_text,
            "notes": notes
        })

    # 4. Calculate Scores
    print("3. Calculating Metrics...")
    
    sem_eval = SemanticEvaluator()
    
    total_rouge = 0
    total_semantic = 0
    structure_hits = 0
    valid_slides = 0
    
    narration_scores = []
    
    results_detail = []

    for slide in slides_data:
        # Skip Title Slide and End Slide heuristics
        if "generate" in slide['title'].lower() and "summary" in slide['title'].lower(): # Title slide
            continue 
        if "thank" in slide['text'].lower() or "questions" in slide['text'].lower():
            continue
            
        valid_slides += 1
        
        # Attempt to align slide to PDF section
        # Normalize slide title 
        slide_key = normalize_heading(slide['title'])
        slide_key_clean = slide['title'].replace("(continued)", "").strip().lower()
        
        # Find best match in PDF
        matched_section_text = ""
        match_type = "None"
        
        # Direct match
        if slide_key in pdf_map:
            matched_section_text = pdf_map[slide_key]
            match_type = "Exact"
        else:
            # Fuzzy match keys
            for p_key, p_text in pdf_map.items():
                if p_key in slide_key_clean or slide_key_clean in p_key:
                    matched_section_text = p_text
                    match_type = "Fuzzy"
                    break
        
        # SCORES
        rouge = 0.0
        semantic = 0.0
        structure_score_slide = 0.0
        
        if matched_section_text:
            rouge = calculate_rouge_l(matched_section_text, slide['text'])
            semantic = sem_eval.compute_similarity(matched_section_text, slide['text'])
            structure_hits += 1
            structure_score_slide = 1.0
        
        # Combine into a final weighted score
        # Since slide generation heavily summarizes text, ROUGE-L will naturally be low.
        # We increase the Semantic weight to 80% to focus on meaning, and lower ROUGE & Structure.
        semantic_norm = max(0.0, semantic)
        weighted_score = (semantic_norm * 0.80) + (rouge * 0.10) + (structure_score_slide * 0.10)
        
        # Adjusted Thresholds for Summarized Academic Content
        if weighted_score >= 0.50:
            classification = "Good"
        elif weighted_score >= 0.30:
            classification = "Average"
        else:
            classification = "Poor"
        
        # Narration
        fk_score = calculate_readability(slide['notes'])
        narration_scores.append(fk_score)
        
        results_detail.append({
            "slide": slide['title'],
            "match_type": match_type,
            "structure_score": structure_score_slide,
            "rouge_l": round(rouge, 4),
            "semantic_sim": round(semantic, 4),
            "weighted_score": round(weighted_score, 4),
            "classification": classification,
            "narration_readability": round(fk_score, 2)
        })

    # Aggregates
    avg_rouge = sum(d['rouge_l'] for d in results_detail) / valid_slides if valid_slides else 0
    avg_semantic = sum(d['semantic_sim'] for d in results_detail) / valid_slides if valid_slides else 0
    structure_score = structure_hits / valid_slides if valid_slides else 0
    avg_weighted_score = sum(d['weighted_score'] for d in results_detail) / valid_slides if valid_slides else 0
    avg_narr_readability = sum(narration_scores) / len(narration_scores) if narration_scores else 0
    
    class_counts = {"Good": 0, "Average": 0, "Poor": 0}
    for d in results_detail:
        class_counts[d['classification']] += 1
    
    # Final Report
    if avg_weighted_score >= 0.50:
        overall_classification = "Good"
    elif avg_weighted_score >= 0.30:
        overall_classification = "Average"
    else:
        overall_classification = "Poor"

    report = {
        "overall_scores": {
            "structure_alignment_score": round(structure_score, 2),
            "semantic_fidelity_score": round(avg_semantic, 2), # 0-1 (Embedding Cosine)
            "lexical_overlap_score_rougel": round(avg_rouge, 4),
            "overall_weighted_score": round(avg_weighted_score, 4),
            "overall_classification": overall_classification,
            "avg_narration_readability_grade": round(avg_narr_readability, 2)
        },
        "slide_classification_summary": class_counts,
        "slide_details": results_detail
    }
    
    print("\n" + "="*40)
    print("       EVALUATION REPORT")
    print("="*40)
    print(f"Structure Alignment: {report['overall_scores']['structure_alignment_score']*100}% of slides matched content.")
    print(f"Content Fidelity (Semantic): {report['overall_scores']['semantic_fidelity_score']:.3f} (Weight: 80%)")
    print(f"Lexical Overlap (ROUGE-L):   {report['overall_scores']['lexical_overlap_score_rougel']:.3f} (Weight: 10%)")
    print(f"Overall Weighted Score:      {report['overall_scores']['overall_weighted_score']:.3f} / 1.0 (Overall: {overall_classification})")
    print(f"Slide Quality Distribution:  {class_counts['Good']} Good, {class_counts['Average']} Average, {class_counts['Poor']} Poor")
    print(f"Narration Grade Level:       {report['overall_scores']['avg_narration_readability_grade']}")
    print("-" * 40)
    
    # Save
    with open(output_json, 'w') as f:
        json.dump(report, f, indent=2)
        
    print(f"Detailed JSON report saved to: {output_json}")

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Evaluate Slide Generation Quality")
    parser.add_argument("pdf", help="Path to original PDF")
    parser.add_argument("pptx", help="Path to generated PPTX")
    args = parser.parse_args()
    
    evaluate_pipeline(args.pdf, args.pptx)
