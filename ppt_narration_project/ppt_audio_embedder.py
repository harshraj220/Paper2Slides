from pptx import Presentation
from pptx.util import Inches
import os


def create_default_icon(path):
    # Create simple 100x100 blue square with 'Audio' text if possible, else just a colored square
    try:
        from PIL import Image, ImageDraw
        img = Image.new('RGB', (100, 100), color = (73, 109, 137))
        d = ImageDraw.Draw(img)
        # We might not have a font easily, so just draw a rectangle or simple shape
        d.rectangle([25, 25, 75, 75], fill=(255, 255, 0))
        img.save(path)
    except ImportError:
        # Fallback if PIL not available (though likely is)
        pass

def embed_audio(input_ppt, audio_dir, output_ppt):
    prs = Presentation(input_ppt)
    
    poster = "audio_icon.png"
    if not os.path.exists(poster):
        create_default_icon(poster)
        
    for idx, slide in enumerate(prs.slides, start=1):
        audio_path = f"{audio_dir}/slide_{idx}.mp3"

        try:
            # Position at bottom left
            left = Inches(0.5)
            top = Inches(6.5)
            width = Inches(1)
            height = Inches(1)
            
            # Check availability again
            use_poster = poster if os.path.exists(poster) else None

            slide.shapes.add_movie(
                audio_path,
                left,
                top,
                width,
                height,
                poster_frame_image=use_poster,
                mime_type="audio/mp3"
            )
        except FileNotFoundError:
            print(f"[WARN] Audio not found for slide {idx}")

    prs.save(output_ppt)
